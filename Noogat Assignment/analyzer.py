import os
import argparse
import io
import re
from PIL import Image
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv
from rich.console import Console
from rich.markdown import Markdown

# --- SETUP AND INITIALIZATION ---
# This section sets up the script's environment.
console = Console()
# Loads the GOOGLE_API_KEY from the .env file for security.
load_dotenv()
api_key = os.getenv("GOOGLE_API_KEY")

if not api_key:
    console.print("[bold red]ERROR: GOOGLE_API_KEY not found.[/bold red]")
    console.print("Please create a .env file and add your Google API Key to it.")
    exit()

# Configures the Gemini client and initializes the AI model.
try:
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel('gemini-1.5-flash')
except Exception as e:
    console.print(f"[bold red]Failed to configure Generative AI: {e}[/bold red]")
    exit()

# --- EXTRACTION FUNCTIONS ---

def extract_content_from_pptx(pptx_path):
    """Extracts all text from a .pptx file, including text from shapes and embedded images."""
    try:
        presentation = Presentation(pptx_path)
    except Exception as e:
        console.print(f"[bold red]Error opening presentation file: {e}[/bold red]")
        return None

    full_text_content = ""
    console.print(f"[cyan]Processing {len(presentation.slides)} slides from .pptx file...[/cyan]")

    # Loop through each slide to extract its content.
    for i, slide in enumerate(presentation.slides):
        slide_number = i + 1
        full_text_content += f"--- Slide {slide_number} ---\n\n"
        
        # 1. Get standard text from shapes like text boxes and tables.
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text_content += shape.text + "\n"

        # 2. Use AI to get text from images embedded in the slide.
        for shape in slide.shapes:
            # Shape type 13 identifies a Picture object.
            if shape.shape_type == 13:
                image = shape.image
                img = Image.open(io.BytesIO(image.blob))
                # Send the image to Gemini for Optical Character Recognition (OCR).
                response = model.generate_content([
                    "Extract all text verbatim from this image. If no text is present, say nothing.",
                    img
                ])
                if response.candidates and not response.candidates[0].finish_reason.name == "SAFETY":
                    full_text_content += f"[Text from image on slide]: {response.text}\n"
    return full_text_content

def extract_content_from_image_folder(folder_path):
    """Extracts text by performing OCR on all images in a specified folder."""
    if not os.path.isdir(folder_path):
        console.print(f"[bold red]Error: Folder not found at '{folder_path}'[/bold red]")
        return None
    
    # Get all image files and sort them naturally (e.g., slide2.png before slide10.png).
    image_files = [f for f in os.listdir(folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp'))]
    image_files.sort(key=lambda f: int(re.sub(r'\D', '', f) or 0))
    
    if not image_files:
        console.print(f"[bold yellow]Warning: No image files found in '{folder_path}'[/bold yellow]")
        return ""

    full_text_content = ""
    console.print(f"[cyan]Processing {len(image_files)} images from folder...[/cyan]")
    
    # Loop through each image file, perform OCR, and add its text to the content string.
    for i, filename in enumerate(image_files):
        slide_number = i + 1
        full_text_content += f"--- Slide {slide_number} ({filename}) ---\n\n"
        
        try:
            img = Image.open(os.path.join(folder_path, filename))
            response = model.generate_content([
                "Extract all text verbatim from this image. If no text is present, say nothing.",
                img
            ])
            if response.candidates and not response.candidates[0].finish_reason.name == "SAFETY":
                full_text_content += response.text + "\n"
        except Exception as e:
            full_text_content += f"[Error processing image {filename}: {e}]\n"

    return full_text_content

# --- ANALYSIS FUNCTION ---

def analyze_content_with_gemini(content):
    """Sends the extracted content to Gemini for inconsistency analysis."""
    # The system prompt is the instruction manual for the AI.
    # It defines its role, task, and the required output format.
    system_prompt = """
    You are a meticulous business analyst. Your task is to review the provided content, extracted slide-by-slide from a presentation, and identify all factual or logical inconsistencies.

    For each inconsistency you find, provide a clear, structured report using the following format:
    **Inconsistency Found:**
    - **Slides Involved:** [e.g., Slide 2 and Slide 5]
    - **Conflicting Information:** [Quote the specific conflicting pieces of data or text]
    - **Analysis:** [Briefly explain why this is an inconsistency]
    ---
    If you find no inconsistencies, your only response should be: "No inconsistencies found."
    """
    
    # The final prompt combines the instructions with the actual slide data.
    prompt = f"{system_prompt}\n\nHere is the presentation content:\n\n{content}"
    
    console.print("[bold cyan]Analyzing content with Gemini... This may take a moment.[/bold cyan]")
    response = model.generate_content(prompt)
    return response.text

# --- SCRIPT EXECUTION ---

def main():
    """Main function to parse arguments and run the analysis."""
    # Set up the command-line interface. The user must supply either --pptx or --image_folder.
    parser = argparse.ArgumentParser(description="Analyze presentation content for inconsistencies.")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--pptx", help="Path to the .pptx file.")
    group.add_argument("--image_folder", help="Path to the folder with slide images.")
    
    args = parser.parse_args()

    # Call the correct extraction function based on the user's input.
    extracted_data = ""
    if args.pptx:
        extracted_data = extract_content_from_pptx(args.pptx)
    elif args.image_folder:
        extracted_data = extract_content_from_image_folder(args.image_folder)

    # If data was extracted successfully, proceed with analysis and reporting.
    if extracted_data:
        final_report = analyze_content_with_gemini(extracted_data)
        
        # Print a nicely formatted report to the terminal using the rich library.
        console.print("\n" + "="*50)
        console.print("         AI Inconsistency Report", style="bold white on blue")
        console.print("="*50 + "\n")
        
        md = Markdown(final_report)
        console.print(md)

# This standard Python construct ensures that the main() function runs only
# when the script is executed directly from the command line.
if __name__ == "__main__":
    main()
